"""
day1_phl101_app_full.py

Comprehensive Streamlit app for Day 1 — PHL 101
Features added on user request:
 - Short-answer auto-grading (keyword-based, configurable)
 - Example Flask-based LLM proxy backend (for secure server-side API calls)
 - PowerPoint export (python-pptx) so you can import to Google Slides
 - All features from previous app: quizzes, journaling, timers, LLM Q&A

Usage:
  - Install requirements: pip install streamlit openai python-pptx flask requests
  - To run Streamlit app: streamlit run day1_phl101_app_full.py
  - To run the LLM proxy server locally (optional):
        python day1_phl101_app_full.py --run-proxy
    The proxy will read OPENAI_API_KEY from environment variables and expose /api/ask

Security note: For production, run the LLM proxy behind HTTPS and authentication. Do NOT embed API keys in the frontend.
"""

import streamlit as st
import io, time, os, argparse, json, threading
from typing import List, Dict

# Optional imports for functionality that may not be present in minimal envs
try:
    import openai
except Exception:
    openai = None
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
except Exception:
    Presentation = None
try:
    from flask import Flask, request, jsonify
    import requests
except Exception:
    Flask = None

# ----------------- App data -----------------
SLIDES = [
    {"title": "Title", "subtitle": "🎓 What is Religion? What is Philosophy?",
     "content": "PHL 101 – Comparative Religions I\nProfessor Xavier Honablue, M.Ed.",
     "notes": "Welcome students warmly. Share your background briefly. Set collaborative tone."},
    {"title": "Course Objectives", "subtitle": "🎯 Our Journey Together",
     "content": "World religions, Indigenous traditions, philosophical approaches, critical thinking, personal reflection.",
     "notes": "Emphasize comparative approach and respect for traditions."},
    {"title": "Icebreaker", "subtitle": "🤔 The Big Questions",
     "content": "Pair Discussion (10 minutes): What is religion? What is philosophy? Where do they overlap?",
     "timer_minutes": 10,
     "notes": "Give students the full 10 minutes; walk the room; collect word cloud ideas."},
    {"title": "Philosophy Meets Religion", "subtitle": "🧠 Philosophy Meets Religion",
     "content": "Definitions, core questions shared by both (ultimate reality, meaning, morality).",
     "notes": "Explain etymology, show overlap."},
    {"title": "Examples", "subtitle": "💡 A Tale of Two Searches",
     "content": "Plato's Cave vs Moses and the Exodus — journeys from ignorance to truth.",
     "notes": "Use the two stories to show shared narrative arc."},
    {"title": "Group Activity", "subtitle": "🎲 Sorting the Big Questions",
     "content": "Sort questions into Philosophy / Religion / Both (15 minutes).",
     "timer_minutes": 15,
     "notes": "Encourage debate; most will end up 'Both'."},
    {"title": "Defining Religion", "subtitle": "🔬 How Scholars Define Religion",
     "content": "Durkheim (social glue), Tylor (belief in spiritual beings), Tillich (ultimate concern).",
     "notes": "Ask which resonates and why."},
    {"title": "Interactive Debate", "subtitle": "💬 Let's Debate!",
     "content": "Team Durkheim, Team Tylor, Team Tillich. Is Buddhism a religion under each definition?",
     "notes": "Moderate debate; present 'challenge question'."},
    {"title": "Wrap-Up & Homework", "subtitle": "🎯 Exit Ticket & Next Steps",
     "content": "Exit ticket: One question about life or religion you hope the class will answer.\nHomework: CrashCourse/TED-Ed videos + 1 paragraph reflection.",
     "notes": "Collect exit tickets & preview Day 2."}
]

QUIZZES = [
    {
        "id": "quiz1",
        "title": "Quick Check: Definitions",
        "questions": [
            {"q": "Durkheim said religion is primarily what?", "choices": ["Belief in gods", "Social glue", "Ultimate concern", "Personal morality"],
             "answer": 1, "explain": "Durkheim emphasized religion's social function — shared rituals and solidarity."},
            {"q": "Tylor's definition focuses on what?", "choices": ["Community", "Ritual", "Belief in spiritual beings", "Ultimate concerns"],
             "answer": 2, "explain": "Tylor defined religion as belief in spiritual beings."}
        ],
        "passing_pct": 60
    }
]

DAY2_PREVIEW = """
Day 2 Preview — Argument Structure
---------------------------------
We'll begin with core tools of argument analysis:
 - Contradiction, Logic, Fallacy, Absurdity (Reductio ad absurdum)
 - Practice identifying premises, conclusions, validity & soundness.
"""

# ----------------- Utilities -----------------

def short_answer_score(student_answer: str, target_keywords: List[str], threshold: float = 0.6) -> Dict:
    """
    Very simple keyword-based short-answer grader.
    - student_answer: raw text from student
    - target_keywords: list of keywords/phrases that indicate coverage
    - threshold: fraction of keywords that must be present to 'pass'

    Returns: {matched: int, total: int, pct: float, passed: bool, matches: list}
    """
    if not student_answer:
        return {"matched": 0, "total": len(target_keywords), "pct": 0.0, "passed": False, "matches": []}
    s = student_answer.lower()
    matches = []
    for kw in target_keywords:
        if kw.lower() in s:
            matches.append(kw)
    matched = len(matches)
    total = max(1, len(target_keywords))
    pct = matched / total
    return {"matched": matched, "total": total, "pct": pct, "passed": pct >= threshold, "matches": matches}

# ----------------- PPTX export -----------------

def generate_pptx_from_slides(slides: List[Dict]) -> bytes:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Install with: pip install python-pptx")
    prs = Presentation()
    # Choose a simple title slide layout for the first and a content layout for others
    for i, s in enumerate(slides):
        if i == 0:
            slide_layout = prs.slide_layouts[0]
        else:
            slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        # Title
        try:
            title = slide.shapes.title
            title.text = s.get('subtitle', s.get('title',''))
        except Exception:
            pass
        # Content body
        try:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            paragraphs = s.get('content','').split('\n')
            for p in paragraphs:
                p_par = tf.add_paragraph()
                p_par.text = p
                p_par.font.size = Pt(14)
        except Exception:
            pass
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()

# ----------------- Minimal Flask proxy (example) -----------------

def create_llm_proxy_app():
    if Flask is None:
        raise RuntimeError("Flask is not installed. Install with: pip install flask")
    app = Flask(__name__)

    @app.route('/api/ask', methods=['POST'])
    def ask():
        data = request.get_json() or {}
        question = data.get('question', '')
        if not question:
            return jsonify({'error': 'question required'}), 400
        api_key = os.environ.get('OPENAI_API_KEY')
        if not api_key:
            return jsonify({'error': 'OPENAI_API_KEY not set on server'}), 500
        if openai is None:
            return jsonify({'error': 'openai package not installed on server'}), 500
        try:
            openai.api_key = api_key
            prompt = [
                {"role": "system", "content": "You are a helpful teaching assistant for an introductory philosophy/religion class. Keep answers concise and reference Day 1 topics."},
                {"role": "user", "content": question}
            ]
            resp = openai.ChatCompletion.create(model="gpt-4o", messages=prompt, max_tokens=400, temperature=0.2)
            answer = resp['choices'][0]['message']['content']
            return jsonify({'answer': answer})
        except Exception as e:
            return jsonify({'error': str(e)}), 500

    return app

# ----------------- Streamlit UI -----------------

def main():
    st.set_page_config(page_title="Day 1 — PHL 101 (Full)", page_icon="🎓", layout="wide")

    # Session defaults
    if 'slide_idx' not in st.session_state:
        st.session_state.slide_idx = 0
    if 'journals' not in st.session_state:
        st.session_state.journals = []
    if 'quiz_scores' not in st.session_state:
        st.session_state.quiz_scores = {}
    if 'llm_history' not in st.session_state:
        st.session_state.llm_history = []

    # Sidebar: config
    with st.sidebar:
        st.title("Day 1 — PHL 101")
        if st.button("Prev"):
            st.session_state.slide_idx = max(0, st.session_state.slide_idx - 1)
        if st.button("Next"):
            st.session_state.slide_idx = min(len(SLIDES)-1, st.session_state.slide_idx + 1)
        st.markdown("---")
        st.subheader("LLM / Proxy")
        openai_key = st.text_input("(Optional) OpenAI API key (direct calls)", type='password')
        proxy_url = st.text_input("(Optional) LLM proxy URL (e.g. https://yourserver/api/ask)")
        use_proxy = st.checkbox("Use proxy instead of direct key", value=False)
        st.markdown("---")
        st.subheader("Export & Utilities")
        if st.button("Export PPTX of slides"):
            try:
                pptx_bytes = generate_pptx_from_slides(SLIDES)
                st.download_button("Download PowerPoint (.pptx)", data=pptx_bytes, file_name="Day1_PHL101.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation")
                st.success("Generated PPTX. You can import the file into Google Slides.")
            except Exception as e:
                st.error(f"Could not generate PPTX: {e}")
        st.markdown("Import the PPTX into Google Slides via File > Import > Upload.")

    # Top header
    st.markdown("# 🎓 Day 1 — What is Religion? What is Philosophy?")
    st.write("---")

    left, right = st.columns([3,1])

    # Left: slide
    slide = SLIDES[st.session_state.slide_idx]
    with left:
        st.subheader(slide.get('subtitle'))
        st.markdown(f"### {slide.get('title')}")
        st.write(slide.get('content'))
        if 'timer_minutes' in slide:
            if st.button(f"Start {slide['timer_minutes']} min timer"):
                st.session_state.timer_seconds = int(slide['timer_minutes'] * 60)
                st.session_state.timer_start = time.time()
        if st.checkbox("Show presenter notes"):
            st.info(slide.get('notes',''))

    # Right: journal, short-answer grading, quizzes, LLM
    with right:
        st.markdown("### Student Journal")
        ch = st.text_area("Describe a challenge you're facing:", key='cj')
        effort = st.text_area("Effort you've made:", key='je')
        if st.button("Save Journal Entry"):
            st.session_state.journals.append({'challenge': ch, 'effort': effort, 'slide': st.session_state.slide_idx, 'ts': time.time()})
            st.success("Saved")

        st.markdown("---")
        st.markdown("### Short-answer auto-grader")
        st.markdown("Give a short-answer prompt using keywords to grade against.")
        sa_prompt = st.text_input("Prompt (e.g., 'Define 'religion' according to Durkheim')")
        sa_answer = st.text_area("Student answer (paste or type)")
        kw_input = st.text_input("Keywords/phrases (comma-separated)", value="social glue, ritual, solidarity")
        threshold = st.slider("Pass threshold (fraction of keywords present)", min_value=0.0, max_value=1.0, value=0.6)
        if st.button("Grade answer"):
            keywords = [k.strip() for k in kw_input.split(',') if k.strip()]
            result = short_answer_score(sa_answer, keywords, threshold)
            st.write(result)
            # show actionable feedback
            if result['passed']:
                st.success(f"Pass — matched {result['matched']} of {result['total']} keywords")
            else:
                st.warning(f"Needs improvement — matched {result['matched']} of {result['total']} keywords. Missing: {set(keywords)-set(result['matches'])}")

        st.markdown("---")
        st.markdown("### Quizzes")
        for q in QUIZZES:
            if st.button(f"Take {q['title']}"):
                st.session_state.current_quiz = q['id']
        if st.session_state.get('current_quiz'):
            q = next(filter(lambda x: x['id']==st.session_state['current_quiz'], QUIZZES))
            st.markdown(f"#### {q['title']}")
            score = 0
            for i, item in enumerate(q['questions']):
                ans = st.radio(item['q'], options=item['choices'], key=f"q_{i}")
                if st.button(f"Submit Q{i+1}", key=f"submit_q_{i}"):
                    chosen_index = item['choices'].index(ans)
                    if chosen_index == item['answer']:
                        st.success("Correct")
                        score += 1
                    else:
                        st.error(f"Incorrect — {item['explain']}")
            if st.button("Finish quiz"):
                pct = int((score / len(q['questions']))*100)
                st.session_state.quiz_scores[q['id']] = {'score': score, 'pct': pct}
                st.success(f"Finished: {score}/{len(q['questions'])} ({pct}%)")

        st.markdown("---")
        st.markdown("### LLM Q&A")
        user_q = st.text_input("Ask the course assistant a question:")
        if st.button("Ask"):
            answer = None
            # prefer proxy if checked
            if st.session_state.get('proxy_url') or st.sidebar.checkbox('Use proxy for this session'):
                url = proxy_url or st.session_state.get('proxy_url')
                if url:
                    try:
                        resp = requests.post(url, json={'question': user_q}, timeout=12)
                        if resp.status_code == 200:
                            answer = resp.json().get('answer')
                        else:
                            answer = f"Proxy error {resp.status_code}: {resp.text}"
                    except Exception as e:
                        answer = f"Proxy request failed: {e}"
                else:
                    answer = "No proxy URL configured. Provide a proxy URL in the sidebar."
            elif openai_key and openai is not None:
                try:
                    openai.api_key = openai_key
                    prompt = [{"role":"system","content":"You are a helpful TA for Day 1 of PHL 101."}, {"role":"user","content": user_q}]
                    resp = openai.ChatCompletion.create(model='gpt-4o', messages=prompt, max_tokens=400, temperature=0.2)
                    answer = resp['choices'][0]['message']['content']
                except Exception as e:
                    answer = f"OpenAI call failed: {e}"
            else:
                answer = "No LLM configured. Provide OpenAI key or a proxy URL in the sidebar."
            st.markdown("**Answer:**")
            st.write(answer)
            st.session_state.llm_history.append({'q': user_q, 'a': answer})

    st.write("---")
    if st.session_state.get('timer_seconds'):
        remaining = st.session_state.timer_seconds - int(time.time() - st.session_state.timer_start)
        if remaining <= 0:
            st.success("Timer finished 🔔")
            st.session_state.pop('timer_seconds', None)
            st.session_state.pop('timer_start', None)
        else:
            mins = remaining // 60
            secs = remaining % 60
            st.info(f"Timer: {mins}:{str(secs).zfill(2)} remaining")

    st.markdown("---")
    st.markdown("### Day 2 Preview")
    st.write(DAY2_PREVIEW)

# ----------------- CLI for running Flask proxy -----------------

def run_proxy_server(port: int = 8765):
    app = create_llm_proxy_app()
    app.run(host='0.0.0.0', port=port)

# ----------------- Entrypoint -----------------
if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('--run-proxy', action='store_true', help='Start the example Flask LLM proxy server')
    parser.add_argument('--proxy-port', type=int, default=8765, help='Port for proxy')
    args = parser.parse_args()

    if args.run_proxy:
        print('Starting LLM proxy server on port', args.proxy_port)
        run_proxy_server(port=args.proxy_port)
    else:
        main()
